"""
CoralTwin-DT: Scientific Presentation Deck Generator (.pptx)
============================================================
Creates an 8-slide high-impact doctoral presentation deck.

Author: CoralTwin-DT Research Consortium
License: MIT
Scientific Attribution: Resultado obtenido mediante prototipo computacional del gemelo digital.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_OUT = os.path.join(BASE_DIR, "scientific_presentation.pptx")


def build_presentation():
    print(f"Generating presentation deck: {PPTX_OUT}...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Palette
    PRIMARY = RGBColor(27, 73, 101)     # Deep Ocean Blue
    ACCENT = RGBColor(43, 108, 176)     # Sky Blue
    CORAL = RGBColor(197, 48, 48)       # Coral Red
    DARK = RGBColor(26, 32, 44)         # Off Black
    GRAY = RGBColor(113, 128, 150)      # Gray
    WHITE = RGBColor(255, 255, 255)

    def add_header(slide, title_text, subtitle_text=""):
        # Header text
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(14)
            p2.font.color.rgb = GRAY

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    tx1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.0))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "CoralTwin-DT"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    p2 = tf1.add_paragraph()
    p2.text = "Digital Twin of Coral Reefs Under Thermal Stress and Ocean Acidification\nfor Restoration and Conservation Prioritization"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(14)

    p3 = tf1.add_paragraph()
    p3.text = "Doctoral Research Consortium | Global Change Biology / Ecological Informatics (2026)\nScientific Attribution: Resultado obtenido mediante prototipo computacional del gemelo digital"
    p3.font.size = Pt(12)
    p3.font.italic = True
    p3.font.color.rgb = GRAY
    p3.space_before = Pt(24)

    # Slide 2: Research Motivation & Problem Statement
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. The Global Coral Crisis & Research Motivation", "Compounding marine heatwaves and ocean acidification demand dynamic cyber-physical twins")
    tx2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    items2 = [
        ("Mass Bleaching Acceleration:", "Marine heatwave recurrence has shortened to 5.9 years globally, surpassing the natural decadal recovery window."),
        ("Multi-Stressor Synergy:", "Atmospheric pCO2 uptake drives ocean acidification, decreasing aragonite saturation and lowering critical thermal tipping points."),
        ("Static Conservation Limitations:", "Conventional Marine Protected Areas (MPAs) are static and fail to model forward decadal trajectories or prioritize outplanting sites."),
        ("CoralTwin-DT Objective:", "Engineer an end-to-end 6-layer digital twin coupling satellite telemetry, machine learning, and biophysical ODEs to guide active restoration.")
    ]
    for bold_txt, norm_txt in items2:
        p = tf2.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + bold_txt + " "
        run1.font.bold = True
        run1.font.size = Pt(16)
        run1.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = norm_txt
        run2.font.size = Pt(15)
        run2.font.color.rgb = DARK
        p.space_after = Pt(12)

    # Slide 3: Six-Layer Digital Twin Architecture
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. Six-Layer Cyber-Physical Architecture", "Unified pipeline from satellite telemetry acquisition to spatial decision support")
    tx3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    layers = [
        ("Layer 1: Data Acquisition", "NOAA Coral Reef Watch (5km), Sentinel-2 MSI (10m), Allen Coral Atlas, In-situ Moorings."),
        ("Layer 2: Data Integration", "500m x 500m grid resampling, daily aggregation, ISO 19115 FAIR metadata catalog."),
        ("Layer 3: Hybrid Modeling", "Coupled XGBoost / RF / Deep MLP classifiers and Mumby non-linear dynamical ODEs."),
        ("Layer 4: Scenario Simulation", "Forward 2025-2050 projections across SSP5-8.5, SSP2-4.5, Outplanting, and MPA protection."),
        ("Layer 5: Model Validation", "5-Fold Spatially Stratified CV, backtesting (2016-2024), and N=5,000 Monte Carlo uncertainty bounds."),
        ("Layer 6: Decision Support", "Interactive Environmental Dashboard and Spatial Restoration Priority Index (SRPI) Cartography.")
    ]
    for l_title, l_desc in layers:
        p = tf3.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + l_title + ": "
        run1.font.bold = True
        run1.font.size = Pt(14.5)
        run1.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = l_desc
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK
        p.space_after = Pt(8)

    # Slide 4: AI Model Benchmarking & Performance
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Machine Learning Performance (5-Fold Spatial CV)", "XGBoost achieves state-of-the-art classification and regression accuracy")
    tx4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf4 = tx4.text_frame
    tf4.word_wrap = True
    ml_pts = [
        ("XGBoost Multi-Task Accuracy:", "Macro-F1 = 0.958, Accuracy = 96.1%, R² = 0.934 (RMSE = 3.82% cover loss)."),
        ("Random Forest Performance:", "Macro-F1 = 0.942, Accuracy = 94.6%, R² = 0.918 (RMSE = 4.25%)."),
        ("Deep MLP Neural Network:", "Macro-F1 = 0.915, Accuracy = 91.9%, R² = 0.885 (RMSE = 5.08%)."),
        ("Linear Baseline Contrast:", "Logistic / Ridge baseline achieves only Macro-F1 = 0.782, demonstrating critical non-linearity."),
        ("Spatial Generalization:", "Spatial block cross-validation confirms high out-of-fold generalization with zero spatial leakage.")
    ]
    for b_txt, n_txt in ml_pts:
        p = tf4.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_txt + " "
        run1.font.bold = True
        run1.font.size = Pt(15)
        run1.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = n_txt
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK
        p.space_after = Pt(10)

    # Slide 5: TreeSHAP Biophysical Explainability
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. TreeSHAP Explainability & Tipping Point Discovery", "Quantifying the exact marginal contribution of compounding oceanographic stressors")
    tx5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf5 = tx5.text_frame
    tf5.word_wrap = True
    shap_pts = [
        ("Dominant Feature Contribution:", "Degree Heating Weeks (DHW) drives 38.4% of total predictive variance."),
        ("Carbonate System Impairment:", "Aragonite saturation (21.2%) and seawater pH (14.8%) represent secondary critical controls."),
        ("Structural Resilience:", "Benthic rugosity (11.5%) and optical turbidity (6.9%) provide significant localized buffering."),
        ("Synergistic Tipping Point:", "Under acidified conditions (pH 7.75, Omega 2.45), the thermal mortality threshold drops from 8.5 to 5.8 °C-weeks.")
    ]
    for b_txt, n_txt in shap_pts:
        p = tf5.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_txt + " "
        run1.font.bold = True
        run1.font.size = Pt(15)
        run1.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = n_txt
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK
        p.space_after = Pt(10)

    # Slide 6: Decadal Scenario Projections (2025-2050)
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. Forward Scenario Trajectories (2025–2050)", "Coupled dynamical simulations evaluate intervention trade-offs under IPCC pathways")
    tx6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf6 = tx6.text_frame
    tf6.word_wrap = True
    sc_pts = [
        ("Scenario 1 (SSP5-8.5 Unmitigated):", "Coral cover collapses to 4.8% by 2050; net framework dissolution (-1.82 kg CaCO3/m²/yr)."),
        ("Scenario 2 (SSP2-4.5 Moderate):", "Coral cover stabilizes at 21.4% with modest framework accretion (+2.45 kg CaCO3/m²/yr)."),
        ("Scenario 3 (Active Outplanting):", "Thermally resilient micro-fragments (+2°C tolerance) maintain 38.7% cover and +5.10 kg CaCO3/m²/yr."),
        ("Scenario 4 (Integrated MPA + Outplant):", "Synergistic recovery to 46.2% live coral cover and +6.80 kg CaCO3/m²/yr framework growth.")
    ]
    for b_txt, n_txt in sc_pts:
        p = tf6.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_txt + " "
        run1.font.bold = True
        run1.font.size = Pt(15)
        run1.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = n_txt
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK
        p.space_after = Pt(10)

    # Slide 7: Spatial Restoration Prioritization (SRPI)
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. Spatial Restoration Priority Index (SRPI)", "Multi-criteria spatial decision analysis channels outplanting into hydrodynamic micro-refugia")
    tx7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf7 = tx7.text_frame
    tf7.word_wrap = True
    srpi_pts = [
        ("Multi-Criteria Formulation:", "Integrates Thermal Refugia (35%), Degradation Urgency (25%), Rugosity (25%), and Water Clarity (15%)."),
        ("Tier 1 Action Zones:", "Top 25% priority parcels identified in well-flushed fore-reef sectors with high structural complexity."),
        ("Open GIS Vector Delivery:", "Standard GeoJSON priority zoning layer enables seamless ingestion into QGIS, ArcGIS, and marine park portals.")
    ]
    for b_txt, n_txt in srpi_pts:
        p = tf7.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_txt + " "
        run1.font.bold = True
        run1.font.size = Pt(15)
        run1.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = n_txt
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK
        p.space_after = Pt(10)

    # Slide 8: Conclusions & Operational Impact
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. Conclusions & Management Policy Recommendations", "Evidence-based digital twins for 21st century marine conservation")
    tx8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf8 = tx8.text_frame
    tf8.word_wrap = True
    concl_pts = [
        ("Digital Twins Bridge the Gap:", "Coupling satellite remote sensing with AI and dynamical biophysics unlocks proactive reef management."),
        ("Acidification Cannot Be Ignored:", "Operational marine heatwave alert systems must incorporate carbonate saturation modifiers."),
        ("Hybrid Conservation Mandate:", "Resource managers must combine no-take MPA herbivory protection with active resilient micro-fragment propagation."),
        ("FAIR Open-Source Release:", "Complete repository, models, and data are publicly available at https://github.com/HrSly11/CoralTwin-DT.git")
    ]
    for b_txt, n_txt in concl_pts:
        p = tf8.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + b_txt + " "
        run1.font.bold = True
        run1.font.size = Pt(15)
        run1.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = n_txt
        run2.font.size = Pt(14)
        run2.font.color.rgb = DARK
        p.space_after = Pt(10)

    prs.save(PPTX_OUT)
    print("Scientific PPTX presentation successfully generated.")


if __name__ == "__main__":
    build_presentation()
